"""
Exports an AnalysisResponse into a client-ready PowerPoint deck.

Slide structure mirrors a real MBB-style deliverable:
  1. Title slide (question + company/industry)
  2. Executive summary / recommendation slide (headline, support, risks)
  3. One slide per issue-tree branch (hypothesis, status, so-what, evidence, quant chart)
  4. Red-team review slide (if present)
  5. Precedents slide (if any were found)

Kept as pure python-pptx (no external template file) so the output has zero
extra dependencies to ship — colors/fonts are set programmatically to match
a simple ink-navy / gold palette consistent with the frontend.
"""
import io
from typing import Optional

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from app.models.schemas import AnalysisResponse, HypothesisStatus

INK = RGBColor(0x0B, 0x1D, 0x33)
GOLD = RGBColor(0xB0, 0x8D, 0x57)
PAPER = RGBColor(0xF7, 0xF5, 0xF0)
SLATE = RGBColor(0x4A, 0x55, 0x68)
SUPPORT = RGBColor(0x2F, 0x6F, 0x4F)
REFUTE = RGBColor(0xB5, 0x53, 0x3C)

STATUS_COLORS = {
    HypothesisStatus.SUPPORTED: SUPPORT,
    HypothesisStatus.PARTIALLY_SUPPORTED: GOLD,
    HypothesisStatus.REFUTED: REFUTE,
    HypothesisStatus.INCONCLUSIVE: SLATE,
    HypothesisStatus.PENDING: SLATE,
}

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _blank_slide(prs: Presentation, bg_color: Optional[RGBColor] = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    if bg_color is not None:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
    return slide


def _add_textbox(slide, left, top, width, height, text, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    return box


def _add_bullets(slide, left, top, width, height, items, size=14, color=INK, bullet_color=GOLD):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = f"—  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        p.space_after = Pt(6)
    return box


def _eyebrow(slide, text, top=Inches(0.4)):
    return _add_textbox(slide, Inches(0.6), top, Inches(8), Inches(0.4), text.upper(), size=12, color=GOLD, bold=True)


def _title_slide(prs: Presentation, analysis: AnalysisResponse):
    slide = _blank_slide(prs, bg_color=INK)
    _add_textbox(slide, Inches(1), Inches(2.7), Inches(11.3), Inches(0.5), "STRATEGY COPILOT", size=14, color=GOLD, bold=True)
    _add_textbox(
        slide, Inches(1), Inches(3.1), Inches(11.3), Inches(2),
        analysis.issue_tree.restated_question, size=32, color=PAPER, bold=True,
    )
    meta_bits = [b for b in [analysis.request.company_name, analysis.request.industry] if b]
    if meta_bits:
        _add_textbox(slide, Inches(1), Inches(5.0), Inches(11.3), Inches(0.5), " · ".join(meta_bits), size=16, color=RGBColor(0xC9, 0xC5, 0xB8))


def _recommendation_slide(prs: Presentation, analysis: AnalysisResponse):
    slide = _blank_slide(prs, bg_color=PAPER)
    _eyebrow(slide, "Recommendation")
    _add_textbox(slide, Inches(0.6), Inches(0.75), Inches(12.1), Inches(1.2), analysis.recommendation.headline, size=26, color=INK, bold=True)
    _add_textbox(slide, Inches(0.6), Inches(1.85), Inches(12.1), Inches(1.3), analysis.recommendation.executive_summary, size=14, color=SLATE)

    _add_textbox(slide, Inches(0.6), Inches(3.3), Inches(5.8), Inches(0.4), "SUPPORTING POINTS", size=12, color=GOLD, bold=True)
    _add_bullets(slide, Inches(0.6), Inches(3.75), Inches(5.8), Inches(3), analysis.recommendation.supporting_points)

    _add_textbox(slide, Inches(6.8), Inches(3.3), Inches(5.8), Inches(0.4), "RISKS & CAVEATS", size=12, color=REFUTE, bold=True)
    _add_bullets(slide, Inches(6.8), Inches(3.75), Inches(5.8), Inches(3), analysis.recommendation.risks_and_caveats, color=SLATE)

    conf_pct = round(analysis.recommendation.confidence * 100)
    _add_textbox(slide, Inches(0.6), Inches(6.9), Inches(4), Inches(0.4), f"Overall confidence: {conf_pct}%", size=12, color=SLATE, bold=True)


def _branch_slide(prs: Presentation, index: int, branch, finding):
    slide = _blank_slide(prs, bg_color=PAPER)
    _eyebrow(slide, f"Branch {index + 1} · {branch.analysis_type.replace('_', ' ').title()}")
    _add_textbox(slide, Inches(0.6), Inches(0.75), Inches(9), Inches(0.8), branch.title, size=24, color=INK, bold=True)

    status_color = STATUS_COLORS.get(finding.status, SLATE) if finding else SLATE
    status_label = finding.status.value.replace("_", " ").title() if finding else "N/A"
    _add_textbox(slide, Inches(10.2), Inches(0.85), Inches(2.5), Inches(0.5), status_label, size=14, color=status_color, bold=True, align=PP_ALIGN.RIGHT)

    _add_textbox(slide, Inches(0.6), Inches(1.55), Inches(12.1), Inches(0.6), f"Hypothesis: {branch.hypothesis}", size=13, color=SLATE)

    if finding:
        _add_textbox(slide, Inches(0.6), Inches(2.3), Inches(3), Inches(0.4), "SO WHAT", size=11, color=GOLD, bold=True)
        _add_textbox(slide, Inches(0.6), Inches(2.7), Inches(6), Inches(1.2), finding.so_what, size=14, color=INK)

        if finding.quant_result and finding.quant_result.method == "tam_sam_som":
            _add_tam_sam_som_chart(slide, finding.quant_result.outputs)
        elif finding.quant_result and finding.quant_result.narrative:
            _add_textbox(slide, Inches(0.6), Inches(4.2), Inches(6), Inches(0.4), "QUANT ANALYSIS", size=11, color=GOLD, bold=True)
            _add_textbox(slide, Inches(0.6), Inches(4.6), Inches(6), Inches(1.5), finding.quant_result.narrative, size=13, color=SLATE)

        if finding.evidence:
            _add_textbox(slide, Inches(7.2), Inches(2.3), Inches(5.5), Inches(0.4), f"EVIDENCE ({len(finding.evidence)})", size=11, color=GOLD, bold=True)
            ev_lines = [f"{e.snippet[:110]}… — {e.source}" for e in finding.evidence[:3]]
            _add_bullets(slide, Inches(7.2), Inches(2.7), Inches(5.5), Inches(3.5), ev_lines, size=11, color=SLATE)

        conf_pct = round(finding.confidence * 100)
        _add_textbox(slide, Inches(0.6), Inches(6.9), Inches(4), Inches(0.4), f"Confidence: {conf_pct}%", size=12, color=SLATE, bold=True)


def _add_tam_sam_som_chart(slide, outputs: dict):
    chart_data = CategoryChartData()
    chart_data.categories = ["TAM", "SAM", "SOM"]
    chart_data.add_series("Value", (outputs.get("tam", 0), outputs.get("sam", 0), outputs.get("som", 0)))
    x, y, cx, cy = Inches(0.6), Inches(4.2), Inches(6), Inches(2.6)
    graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data)
    chart = graphic_frame.chart
    chart.has_legend = False
    plot = chart.plots[0]
    plot.series[0].format.fill.solid()
    plot.series[0].format.fill.fore_color.rgb = GOLD


def _red_team_slide(prs: Presentation, analysis: AnalysisResponse):
    if not analysis.red_team:
        return
    rt = analysis.red_team
    slide = _blank_slide(prs, bg_color=INK)
    _eyebrow(slide, "Red-team review", top=Inches(0.5))
    verdict_color = {"holds": SUPPORT, "weakened": GOLD, "reversed": REFUTE}.get(rt.verdict.value, GOLD)
    _add_textbox(slide, Inches(0.6), Inches(0.95), Inches(8), Inches(0.5), f"Verdict: {rt.verdict.value.title()}", size=18, color=verdict_color, bold=True)

    _add_textbox(slide, Inches(0.6), Inches(1.7), Inches(4), Inches(0.4), "STRONGEST OBJECTION", size=12, color=GOLD, bold=True)
    _add_textbox(slide, Inches(0.6), Inches(2.15), Inches(12), Inches(1), rt.strongest_objection, size=15, color=PAPER)

    _add_textbox(slide, Inches(0.6), Inches(3.3), Inches(5.8), Inches(0.4), "ASSUMPTIONS CHALLENGED", size=12, color=GOLD, bold=True)
    _add_bullets(slide, Inches(0.6), Inches(3.75), Inches(5.8), Inches(2.5), rt.challenged_assumptions, color=PAPER)

    _add_textbox(slide, Inches(6.8), Inches(3.3), Inches(5.8), Inches(0.4), "UNRESOLVED RISKS", size=12, color=GOLD, bold=True)
    _add_bullets(slide, Inches(6.8), Inches(3.75), Inches(5.8), Inches(2.5), rt.unresolved_risks, color=PAPER)

    _add_textbox(
        slide, Inches(0.6), Inches(6.9), Inches(6), Inches(0.4),
        f"Adjusted confidence after review: {round(rt.adjusted_confidence * 100)}%", size=12, color=RGBColor(0xC9, 0xC5, 0xB8), bold=True,
    )


def _precedents_slide(prs: Presentation, analysis: AnalysisResponse):
    if not analysis.precedents:
        return
    slide = _blank_slide(prs, bg_color=PAPER)
    _eyebrow(slide, "Comparable precedents")
    _add_textbox(slide, Inches(0.6), Inches(0.85), Inches(11), Inches(0.6), "What happened when others tried something similar", size=22, color=INK, bold=True)

    top = Inches(1.8)
    for p in analysis.precedents[:3]:
        _add_textbox(slide, Inches(0.6), top, Inches(11.8), Inches(0.4), p.situation, size=15, color=INK, bold=True)
        _add_textbox(slide, Inches(0.6), top + Inches(0.4), Inches(11.8), Inches(0.4), f"Outcome: {p.outcome}", size=13, color=SLATE)
        _add_textbox(slide, Inches(0.6), top + Inches(0.8), Inches(11.8), Inches(0.5), f"Implication: {p.implication}", size=13, color=GOLD, bold=True)
        top += Inches(1.6)


def build_deck(analysis: AnalysisResponse) -> bytes:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    _title_slide(prs, analysis)
    _recommendation_slide(prs, analysis)

    finding_by_branch = {f.branch_id: f for f in analysis.findings}
    for i, branch in enumerate(analysis.issue_tree.branches):
        _branch_slide(prs, i, branch, finding_by_branch.get(branch.id))

    _red_team_slide(prs, analysis)
    _precedents_slide(prs, analysis)

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.read()
